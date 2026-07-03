"""Render a :class:`~quantized.calc.report.ReportSheet` to office / markup files.

ORIGIN_GAP_PLAN #37 (docx/pptx) + #38 (LaTeX) + #39 (HTML). Every renderer
walks the SAME report schema (title -> sections -> typed blocks) with no
per-block special cases beyond the four block kinds — that is the #36
acceptance criterion made real.

LaTeX and HTML are pure-Python and always available. Word (.docx) and
PowerPoint (.pptx) need the MIT libraries ``python-docx`` / ``python-pptx``;
those imports are guarded so the module (and CI) work without them — a missing
library raises a clear ``ReportExportError`` instead of an ImportError at
module load.

Pure ``io`` layer — no fastapi/pydantic imports (enforced by
test_repo_integrity). ``value ± error`` formatting rounds the value to the
precision implied by the uncertainty (2 significant figures on the error).
"""

from __future__ import annotations

import base64 as _base64
import binascii as _binascii
import html as _html
import io as _io
import math
from collections.abc import Mapping
from typing import Any

# Raster image MIME types Office can embed (SVG/other -> placeholder text).
_EMBEDDABLE_IMAGE_MIMES = ("image/png", "image/jpeg", "image/jpg", "image/gif", "image/bmp")


def _decode_raster(image: Mapping[str, str] | None) -> bytes | None:
    """Return decoded image bytes iff it's an Office-embeddable raster type."""
    if not image or image.get("mime") not in _EMBEDDABLE_IMAGE_MIMES:
        return None
    try:
        return _base64.b64decode(image["data"], validate=True)
    except (_binascii.Error, ValueError, KeyError):
        return None

__all__ = [
    "FORMATS",
    "ReportExportError",
    "format_value_error",
    "render_report",
    "to_html",
    "to_latex",
]

FORMATS = ("latex", "html", "docx", "pptx")


class ReportExportError(RuntimeError):
    """Raised for an unknown format or a missing optional export library."""


# ── number formatting ─────────────────────────────────────────────────────
def _fmt_num(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return str(value)
    if isinstance(value, int):
        return str(value)
    if isinstance(value, float):
        if not math.isfinite(value):
            return ""
        return f"{value:.6g}"
    return str(value)


def format_value_error(value: Any, error: Any = None, *, sig: int = 2) -> str:
    """Format ``value ± error`` with the value rounded to the error's precision.

    With no (finite, non-zero) error, falls back to a plain 6-significant-figure
    number. ``sig`` is the number of significant figures kept on the error.
    """
    if value is None:
        return ""
    v = float(value)
    if error is None or not math.isfinite(float(error)) or float(error) == 0.0:
        return _fmt_num(value)
    e = abs(float(error))  # uncertainty is magnitude-only, sign is meaningless
    exp = math.floor(math.log10(e))
    ndp = sig - 1 - exp  # decimal places (may be negative for large errors)
    v_r, e_r = round(v, ndp), round(e, ndp)
    dp = max(0, ndp)
    return f"{v_r:.{dp}f} ± {e_r:.{dp}f}"


# ── block-walking helpers (shared by every renderer) ──────────────────────
def _params_rows(block: Mapping[str, Any]) -> tuple[list[str], list[list[str]]]:
    """(header, rows) for a params block, with value ± error merged."""
    has_unit = any(p.get("unit") for p in block["params"])
    header = ["Parameter", "Value", *(["Unit"] if has_unit else [])]
    rows = []
    for p in block["params"]:
        cells = [str(p["name"]), format_value_error(p.get("value"), p.get("error"))]
        if has_unit:
            cells.append(str(p.get("unit", "")))
        rows.append(cells)
    return header, rows


def _table_rows(block: Mapping[str, Any]) -> tuple[list[str], list[list[str]]]:
    header = [str(c) for c in block["columns"]]
    rows = [[_fmt_num(c) for c in row] for row in block["rows"]]
    return header, rows


# ── LaTeX (booktabs) ──────────────────────────────────────────────────────
# LaTeX special chars + the science glyphs the emitters emit (so the output
# compiles under plain pdfLaTeX, no inputenc/unicode-engine required).
_LATEX_REPL = {
    "&": r"\&", "%": r"\%", "$": r"\$", "#": r"\#", "_": r"\_",
    "{": r"\{", "}": r"\}", "~": r"\textasciitilde{}", "^": r"\textasciicircum{}",
    "±": r"$\pm$", "×": r"$\times$", "·": r"$\cdot$", "²": r"$^2$", "³": r"$^3$",
    "χ": r"$\chi$", "η": r"$\eta$", "α": r"$\alpha$", "β": r"$\beta$",
    "γ": r"$\gamma$", "σ": r"$\sigma$", "λ": r"$\lambda$", "θ": r"$\theta$",
    "ω": r"$\omega$", "μ": r"$\mu$", "π": r"$\pi$", "τ": r"$\tau$",
    "Δ": r"$\Delta$", "Ω": r"$\Omega$", "Å": r"\AA{}", "°": r"$^\circ$",
    "√": r"$\surd$", "∞": r"$\infty$",
}


def _latex_escape(text: str) -> str:
    return "".join(_LATEX_REPL.get(ch, ch) for ch in text)


def _latex_table(header: list[str], rows: list[list[str]], caption: str | None) -> list[str]:
    ncol = len(header)
    align = "l" + "r" * (ncol - 1) if ncol > 1 else "l"
    out = [r"\begin{table}[h]", r"  \centering"]
    if caption:
        out.append(rf"  \caption{{{_latex_escape(caption)}}}")
    out.append(rf"  \begin{{tabular}}{{{align}}}")
    out.append(r"    \toprule")
    out.append("    " + " & ".join(_latex_escape(h) for h in header) + r" \\")
    out.append(r"    \midrule")
    for row in rows:
        out.append("    " + " & ".join(_latex_escape(str(c)) for c in row) + r" \\")
    out.append(r"    \bottomrule")
    out.append(r"  \end{tabular}")
    out.append(r"\end{table}")
    return out


def to_latex(report: Mapping[str, Any]) -> str:
    """Booktabs LaTeX for the report's tables (params + stats), text as prose."""
    lines = [rf"% Report: {_latex_escape(str(report.get('title', '')))}",
             r"% Requires \usepackage{booktabs}", ""]
    for sec in report.get("sections", []):
        lines.append(rf"\subsection*{{{_latex_escape(str(sec.get('title', '')))}}}")
        for block in sec.get("blocks", []):
            btype = block.get("type")
            if btype == "text":
                lines.append(_latex_escape(block["text"]) + "\n")
            elif btype == "params":
                header, rows = _params_rows(block)
                lines += _latex_table(header, rows, block.get("caption"))
            elif btype == "table":
                header, rows = _table_rows(block)
                lines += _latex_table(header, rows, block.get("caption"))
            elif btype == "figure":
                lines.append(rf"% [figure: {_latex_escape(str(block.get('name', '')))}]")
        lines.append("")
    return "\n".join(lines).rstrip() + "\n"


# ── HTML (self-contained) ─────────────────────────────────────────────────
_HTML_CSS = (
    "body{font-family:system-ui,sans-serif;max-width:52rem;margin:2rem auto;"
    "padding:0 1rem;color:#1a1a1a}h1{font-size:1.5rem}h2{font-size:1.15rem;"
    "border-bottom:1px solid #ddd;padding-bottom:.2rem}table{border-collapse:"
    "collapse;margin:.6rem 0}th,td{border:1px solid #ccc;padding:.25rem .6rem;"
    "text-align:right}th:first-child,td:first-child{text-align:left}"
    "caption{caption-side:top;font-style:italic;text-align:left;color:#555}"
    "figure{color:#777;font-style:italic}"
)


def _html_table(header: list[str], rows: list[list[str]], caption: str | None) -> str:
    parts = ["<table>"]
    if caption:
        parts.append(f"<caption>{_html.escape(caption)}</caption>")
    parts.append("<thead><tr>" + "".join(f"<th>{_html.escape(h)}</th>" for h in header)
                 + "</tr></thead><tbody>")
    for row in rows:
        parts.append("<tr>" + "".join(f"<td>{_html.escape(str(c))}</td>" for c in row) + "</tr>")
    parts.append("</tbody></table>")
    return "".join(parts)


def to_html(report: Mapping[str, Any]) -> str:
    """A self-contained HTML page for the report (#39)."""
    title = _html.escape(str(report.get("title", "Report")))
    body = [f"<h1>{title}</h1>"]
    refs = report.get("source_refs", [])
    if refs:
        names = ", ".join(_html.escape(str(r.get("name") or r.get("id"))) for r in refs)
        body.append(f"<p><small>Sources: {names}</small></p>")
    for sec in report.get("sections", []):
        body.append(f"<h2>{_html.escape(str(sec.get('title', '')))}</h2>")
        for block in sec.get("blocks", []):
            btype = block.get("type")
            if btype == "text":
                body.append(f"<p>{_html.escape(block['text'])}</p>")
            elif btype == "params":
                header, rows = _params_rows(block)
                body.append(_html_table(header, rows, block.get("caption")))
            elif btype == "table":
                header, rows = _table_rows(block)
                body.append(_html_table(header, rows, block.get("caption")))
            elif btype == "figure":
                cap = block.get("caption") or block.get("name", "")
                img = block.get("image")
                if img:
                    src = f"data:{img['mime']};base64,{img['data']}"
                    body.append(f'<figure><img src="{src}" alt="{_html.escape(str(cap))}"'
                                f' style="max-width:100%"><figcaption>'
                                f"{_html.escape(str(cap))}</figcaption></figure>")
                else:
                    body.append(f"<figure>[figure: {_html.escape(str(cap))}]</figure>")
    return (f"<!doctype html><html><head><meta charset='utf-8'><title>{title}</title>"
            f"<style>{_HTML_CSS}</style></head><body>{''.join(body)}</body></html>")


# ── Word / PowerPoint (guarded optional deps) ─────────────────────────────
def _to_docx(report: Mapping[str, Any]) -> bytes:
    try:
        from docx import Document  # python-docx (MIT)
    except ImportError as exc:  # pragma: no cover - exercised only without the dep
        raise ReportExportError(
            "Word export needs 'python-docx' (pip install quantized[office])"
        ) from exc

    doc = Document()
    doc.add_heading(str(report.get("title", "Report")), level=0)
    for sec in report.get("sections", []):
        doc.add_heading(str(sec.get("title", "")), level=1)
        for block in sec.get("blocks", []):
            btype = block.get("type")
            if btype == "text":
                doc.add_paragraph(block["text"])
            elif btype in ("params", "table"):
                header, rows = (_params_rows if btype == "params" else _table_rows)(block)
                cap = block.get("caption")
                if cap:
                    doc.add_paragraph().add_run(str(cap)).italic = True
                t = doc.add_table(rows=1, cols=len(header))
                t.style = "Light Grid Accent 1"
                for j, h in enumerate(header):
                    t.rows[0].cells[j].text = h
                for row in rows:
                    cells = t.add_row().cells
                    for j, c in enumerate(row):
                        cells[j].text = str(c)
            elif btype == "figure":
                raster = _decode_raster(block.get("image"))
                if raster is not None:
                    from docx.shared import Inches
                    doc.add_picture(_io.BytesIO(raster), width=Inches(6))
                    if block.get("caption"):
                        doc.add_paragraph().add_run(str(block["caption"])).italic = True
                else:
                    doc.add_paragraph(f"[figure: {block.get('name', '')}]")
    buf = _io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _to_pptx(report: Mapping[str, Any]) -> bytes:
    try:
        from pptx import Presentation  # python-pptx (MIT)
        from pptx.util import Inches, Pt
    except ImportError as exc:  # pragma: no cover - exercised only without the dep
        raise ReportExportError(
            "PowerPoint export needs 'python-pptx' (pip install quantized[office])"
        ) from exc

    prs = Presentation()
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[5]
    first = prs.slides.add_slide(title_layout)
    first.shapes.title.text = str(report.get("title", "Report"))
    for sec in report.get("sections", []):
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        box.text_frame.text = str(sec.get("title", ""))
        box.text_frame.paragraphs[0].font.size = Pt(28)
        top = 1.3
        for block in sec.get("blocks", []):
            btype = block.get("type")
            if btype == "text":
                tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.6))
                tb.text_frame.text = block["text"]
                tb.text_frame.word_wrap = True
                top += 0.7
            elif btype in ("params", "table"):
                header, rows = (_params_rows if btype == "params" else _table_rows)(block)
                nrows, ncols = len(rows) + 1, len(header)
                height = min(0.35 * nrows, 5.0)
                gt = slide.shapes.add_table(
                    nrows, ncols, Inches(0.5), Inches(top), Inches(9), Inches(height)
                ).table
                for j, h in enumerate(header):
                    gt.cell(0, j).text = h
                for i, row in enumerate(rows, start=1):
                    for j, c in enumerate(row):
                        gt.cell(i, j).text = str(c)
                top += height + 0.3
            elif btype == "figure":
                raster = _decode_raster(block.get("image"))
                if raster is not None:
                    slide.shapes.add_picture(
                        _io.BytesIO(raster), Inches(0.5), Inches(top), width=Inches(6)
                    )
                    top += 4.0
                else:
                    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.5))
                    tb.text_frame.text = f"[figure: {block.get('name', '')}]"
                    top += 0.6
    buf = _io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── dispatch ──────────────────────────────────────────────────────────────
def render_report(report: Mapping[str, Any], fmt: str) -> tuple[bytes, str, bool]:
    """Render ``report`` to ``fmt``; return ``(data, mime, is_text)``.

    ``is_text`` is True for latex/html (utf-8 text), False for docx/pptx
    (binary — the route base64-encodes these). Unknown or unavailable formats
    raise :class:`ReportExportError`.
    """
    if fmt == "latex":
        return to_latex(report).encode("utf-8"), "text/x-tex", True
    if fmt == "html":
        return to_html(report).encode("utf-8"), "text/html", True
    if fmt == "docx":
        return (_to_docx(report),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document", False)
    if fmt == "pptx":
        return (_to_pptx(report),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation", False)
    raise ReportExportError(f"unknown report format {fmt!r}; expected one of {FORMATS}")
